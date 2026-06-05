"""
MSN 每日新闻简报生成器
使用 RSS + Bing 新闻搜索获取 AI 和汽车行业头条，整理成 PPT 报告。
执行后将状态写入 task_status.json，供监控面板查看。
"""

import os
import re
import json
import time
import logging
import xml.etree.ElementTree as ET
from datetime import datetime, date, timedelta
from urllib.request import Request, urlopen
from urllib.parse import quote_plus
from typing import List, Dict

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("news_report.log", encoding="utf-8"),
    ],
)
log = logging.getLogger(__name__)

# ======================== 配置 ========================
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "reports")
STATUS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "task_status.json")
MAX_NEWS_PER_CATEGORY = 10
NEXT_RUN_TIME = "每天 07:03"

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                  "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Accept-Language": "zh-CN,zh;q=0.9",
}

RSS_SOURCES = [
    "https://www.ithome.com/rss/",
    "https://36kr.com/feed",
    "https://rss.sina.com.cn/news/marshal/allnews/allnews.xml",
]

BING_SEARCHES = [
    ("ai", "人工智能 OR AI OR 大模型 OR 机器学习"),
    ("ai_gpt", "GPT OR ChatGPT OR Claude OR LLM OR AIGC"),
    ("auto", "汽车 OR 电动车 OR 新能源汽车 OR 特斯拉 OR 比亚迪"),
    ("auto_ev", "新能源汽车 OR 智能驾驶 OR 续航 OR 智驾"),
]


# ======================== 状态管理 ========================
def load_status() -> Dict:
    """Load task status from file."""
    if os.path.exists(STATUS_FILE):
        with open(STATUS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {
        "task_name": "每日新闻简报生成",
        "schedule": "每天 07:03",
        "runs": [],
    }


def save_status(status: Dict):
    """Save task status to file."""
    with open(STATUS_FILE, "w", encoding="utf-8") as f:
        json.dump(status, f, ensure_ascii=False, indent=2)


def record_run(success: bool, ai_count: int, auto_count: int, top_count: int,
               duration: float, ai_items: List[Dict] = None,
               auto_items: List[Dict] = None, top_items: List[Dict] = None,
               error_msg: str = "", ppt_path: str = "", json_path: str = ""):
    """Record a task run result with news items."""
    status = load_status()
    run_record = {
        "timestamp": datetime.now().isoformat(),
        "status": "success" if success else "failed",
        "ai_count": ai_count,
        "auto_count": auto_count,
        "top_count": top_count,
        "total": ai_count + auto_count + top_count,
        "duration": round(duration, 1),
        "error": error_msg,
        "ppt_path": ppt_path,
        "json_path": json_path,
        "ai_items": ai_items or [],
        "auto_items": auto_items or [],
        "top_items": top_items or [],
    }
    status["runs"].append(run_record)
    if len(status["runs"]) > 30:
        status["runs"] = status["runs"][-30:]
    status["last_run"] = run_record
    save_status(status)


def calculate_next_run() -> str:
    """Calculate next scheduled run time."""
    now = datetime.now()
    tomorrow = now + timedelta(days=1)
    next_run = tomorrow.replace(hour=7, minute=3, second=0, microsecond=0)
    if now.hour < 7 or (now.hour == 7 and now.minute < 3):
        next_run = now.replace(hour=7, minute=3, second=0, microsecond=0)
    return next_run.strftime("%Y-%m-%d 07:03")


# ======================== 新闻抓取 ========================
def fetch_url(url: str, timeout: int = 15) -> str:
    """Fetch URL and return text content."""
    try:
        req = Request(url, headers=HEADERS)
        resp = urlopen(req, timeout=timeout)
        return resp.read().decode("utf-8", errors="replace")
    except Exception as e:
        log.warning("Failed to fetch %s: %s", url, e)
        return ""


def parse_rss(xml_content: str, source_name: str) -> List[Dict]:
    """Parse RSS/Atom XML and extract news items."""
    items = []
    try:
        root = ET.fromstring(xml_content)
    except ET.ParseError:
        log.warning("Failed to parse RSS from %s", source_name)
        return items

    for item in root.iter("item"):
        title_el = item.find("title")
        link_el = item.find("link")
        source_el = item.find("source")
        pub_el = item.find("pubDate")

        title = title_el.text.strip() if title_el is not None and title_el.text else ""
        link = link_el.text.strip() if link_el is not None and link_el.text else ""
        src = source_el.text.strip() if source_el is not None and source_el.text else source_name

        if len(title) >= 6:
            items.append({
                "title": re.sub(r"<[^>]+>", "", title)[:120],
                "url": link,
                "source": src or source_name,
                "date": pub_el.text.strip() if pub_el is not None and pub_el.text else "",
            })

    for entry in root.iter("entry"):
        title_el = entry.find(".//{http://www.w3.org/2005/Atom}title")
        link_el = entry.find(".//{http://www.w3.org/2005/Atom}link")
        updated_el = entry.find(".//{http://www.w3.org/2005/Atom}updated")

        if title_el is not None and title_el.text:
            title = title_el.text.strip()
            link = link_el.get("href", "") if link_el is not None else ""

            if len(title) >= 6:
                items.append({
                    "title": re.sub(r"<[^>]+>", "", title)[:120],
                    "url": link,
                    "source": source_name,
                    "date": updated_el.text.strip() if updated_el is not None and updated_el.text else "",
                })

    return items


def search_bing_news(query: str, count: int = 10) -> List[Dict]:
    """Search Bing News via the web interface and extract results."""
    items = []
    encoded_query = quote_plus(query)
    url = f"https://www.bing.com/news/search?q={encoded_query}&setlang=zh-Hans&cc=cn&count={count}"

    html = fetch_url(url, timeout=15)
    if not html:
        return items

    seen_titles = set()

    for match in re.finditer(
        r'<h2[^>]*>\s*<a[^>]*href=["\']([^"\']*)["\'][^>]*>(.*?)</a>\s*</h2>',
        html, re.DOTALL,
    ):
        link = match.group(1)
        title = re.sub(r"<[^>]+>", "", match.group(2)).strip()
        if len(title) >= 6 and title not in seen_titles:
            seen_titles.add(title)
            items.append({
                "title": title[:120], "url": link, "source": "Bing News", "date": "",
            })

    for match in re.finditer(
        r'data-linkurl=["\']([^"\']*)["\'][^>]*data-newsname=["\']([^"\']*)["\']',
        html,
    ):
        link = match.group(1)
        title = re.sub(r"&[^;]*;", "", match.group(2)).strip()
        if len(title) >= 6 and title not in seen_titles:
            seen_titles.add(title)
            items.append({
                "title": title[:120], "url": link, "source": "Bing News", "date": "",
            })

    for match in re.finditer(
        r'<a[^>]*href=["\']([^"\']*)["\'][^>]*class=["\'][^"\']*title[^"\']*["\'][^>]*>(.*?)</a>',
        html, re.DOTALL,
    ):
        link = match.group(1)
        title = re.sub(r"<[^>]+>", "", match.group(2)).strip()
        if len(title) >= 8 and title[:4] not in ("http", "HTTP") and title not in seen_titles:
            seen_titles.add(title)
            items.append({
                "title": title[:120], "url": link, "source": "Bing News", "date": "",
            })

    return items


def fetch_all_news() -> Dict:
    """Fetch news from all sources and classify."""
    all_items = []

    for rss_url in RSS_SOURCES:
        source_name = rss_url.split("/")[-1].split(".")[0]
        log.info("Fetching RSS: %s", rss_url)
        xml_content = fetch_url(rss_url, timeout=15)
        if xml_content:
            items = parse_rss(xml_content, source_name)
            log.info("  Got %d items from %s", len(items), source_name)
            all_items.extend(items)

    for search_key, query in BING_SEARCHES:
        log.info("Searching Bing News: %s", query)
        items = search_bing_news(query, count=10)
        log.info("  Got %d items for %s", len(items), search_key)
        all_items.extend(items)

    seen = set()
    unique_items = []
    for item in all_items:
        if item["title"] not in seen:
            seen.add(item["title"])
            unique_items.append(item)

    ai_keywords = [
        "人工智能", "AI", "大模型", "机器学习", "深度学习",
        "GPT", "ChatGPT", "Claude", "生成式AI", "AIGC",
        "智能驾驶", "自动驾驶", "LLM", "OpenAI",
    ]
    car_keywords = [
        "汽车", "电动车", "新能源汽车", "电动汽车", "比亚迪",
        "特斯拉", "Tesla", "小米汽车", "蔚来", "理想", "小鹏",
        "造车", "销量", "续航", "智驾", "新能源",
        "电池", "充电", "换电", "Robotaxi",
    ]

    ai_news, auto_news, top_news = [], [], []
    for item in unique_items:
        title = item["title"]
        is_ai = any(kw in title for kw in ai_keywords)
        is_auto = any(kw in title for kw in car_keywords)
        if is_ai:
            ai_news.append(item)
        elif is_auto:
            auto_news.append(item)
        else:
            top_news.append(item)

    ai_news = ai_news[:MAX_NEWS_PER_CATEGORY]
    auto_news = auto_news[:MAX_NEWS_PER_CATEGORY]
    top_news = top_news[:MAX_NEWS_PER_CATEGORY]

    log.info("Final: AI=%d, Auto=%d, General=%d",
             len(ai_news), len(auto_news), len(top_news))

    return {
        "date": date.today().isoformat(),
        "ai": ai_news,
        "auto": auto_news,
        "top": top_news,
    }


# ======================== PPT 生成 ========================
def generate_ppt(summary: Dict, output_path: str):
    """Generate a professional PPTX report."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    COLORS = {
        "dark_blue": RGBColor(0x1B, 0x3A, 0x5F),
        "accent_blue": RGBColor(0x00, 0x78, 0xD4),
        "ai_purple": RGBColor(0x7B, 0x2F, 0xBE),
        "auto_green": RGBColor(0x00, 0xA6, 0x6E),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_gray": RGBColor(0xF0, 0xF2, 0xF5),
        "text_dark": RGBColor(0x33, 0x33, 0x33),
        "text_light": RGBColor(0x66, 0x66, 0x66),
        "medium_gray": RGBColor(0xAA, 0xBB, 0xCC),
    }
    FONT = "Microsoft YaHei"

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text(slide, left, top, width, height, text, size=18,
                 color=None, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = FONT
        if color:
            p.font.color.rgb = color
        p.alignment = align
        return txBox

    def add_accent_bar(slide, left, top, width, color):
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["dark_blue"])
    add_text(slide, 1.5, 3.0, 10, 1.2,
             "每日新闻简报", 44, COLORS["white"], True, PP_ALIGN.CENTER)
    add_text(slide, 1.5, 4.5, 10, 0.8,
             "AI · 汽车 · 头条", 24, COLORS["accent_blue"], False, PP_ALIGN.CENTER)
    add_text(slide, 1.5, 5.5, 10, 0.6,
             f"生成日期：{summary['date']}", 16, COLORS["medium_gray"], False, PP_ALIGN.CENTER)

    ai_count = len(summary["ai"])
    auto_count = len(summary["auto"])
    top_count = len(summary["top"])
    total = ai_count + auto_count + top_count

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["light_gray"])
    add_text(slide, 0.8, 0.4, 8, 0.8, "今日概览", 32, COLORS["dark_blue"], True)
    add_text(slide, 1.2, 1.8, 8, 0.6, f"共收集 {total} 条新闻", 22, COLORS["text_dark"])
    add_text(slide, 1.2, 2.6, 8, 0.6, f"AI 领域 {ai_count} 条", 22, COLORS["text_dark"])
    add_text(slide, 1.2, 3.4, 8, 0.6, f"汽车行业 {auto_count} 条", 22, COLORS["text_dark"])

    stats = [
        ("AI 资讯", ai_count, COLORS["ai_purple"]),
        ("汽车资讯", auto_count, COLORS["auto_green"]),
        ("头条精选", top_count, COLORS["accent_blue"]),
    ]
    for i, (label, count, color) in enumerate(stats):
        x = 1.5 + i * 3.5
        add_text(slide, x, 5.0, 3, 1.2, str(count), 44, color, True, PP_ALIGN.CENTER)
        add_text(slide, x, 6.2, 3, 0.5, label, 16, COLORS["text_light"], False, PP_ALIGN.CENTER)

    def add_news_slide(title: str, items: List[Dict], accent_color: RGBColor):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, COLORS["light_gray"])
        add_text(slide, 0.8, 0.4, 8, 0.8, title, 32, COLORS["dark_blue"], True)
        add_accent_bar(slide, 0.8, 1.15, 2, accent_color)

        y = 1.7
        for i, item in enumerate(items[:MAX_NEWS_PER_CATEGORY], 1):
            title_text = item["title"]
            source_text = item.get("source", "")
            if len(title_text) > 55:
                title_text = title_text[:53] + "..."

            badge = slide.shapes.add_shape(
                5, Inches(1.0), Inches(y), Inches(0.45), Inches(0.45),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent_color
            badge.line.fill.background()
            bf = badge.text_frame
            bf.paragraphs[0].text = str(i)
            bf.paragraphs[0].font.size = Pt(14)
            bf.paragraphs[0].font.color.rgb = COLORS["white"]
            bf.paragraphs[0].font.bold = True
            bf.paragraphs[0].alignment = PP_ALIGN.CENTER
            bf.paragraphs[0].font.name = FONT

            add_text(slide, 1.7, y - 0.05, 10, 0.4, title_text, 17, COLORS["text_dark"])

            if source_text:
                add_text(slide, 1.7, y + 0.35, 10, 0.3,
                         f"[{source_text}]", 11, COLORS["text_light"])

            y += 0.75
            if y > 8.2:
                break

    if summary["ai"]:
        add_news_slide("人工智能资讯", summary["ai"], COLORS["ai_purple"])
    if summary["auto"]:
        add_news_slide("汽车行业资讯", summary["auto"], COLORS["auto_green"])
    if summary["top"]:
        add_news_slide("头条精选", summary["top"], COLORS["accent_blue"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLORS["dark_blue"])
    add_text(slide, 1.5, 2.5, 10, 1.0, "感谢您的阅读", 36, COLORS["white"], True, PP_ALIGN.CENTER)
    add_text(slide, 1.5, 4.0, 10, 0.8,
             "数据来源：Bing News · RSS", 18, COLORS["medium_gray"], False, PP_ALIGN.CENTER)
    add_text(slide, 1.5, 4.8, 10, 0.6,
             f"报告生成于 {datetime.now().strftime('%Y-%m-%d %H:%M')}",
             14, RGBColor(0x88, 0x99, 0xAA), False, PP_ALIGN.CENTER)

    prs.save(output_path)
    log.info("PPT saved to: %s", output_path)


# ======================== 主流程 ========================
def main():
    start_time = time.time()
    log.info("=" * 50)
    log.info("每日新闻简报生成器")
    log.info("=" * 50)

    try:
        news_data = fetch_all_news()

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filename = f"每日新闻简报_{news_data['date']}.pptx"
        output_path = os.path.join(OUTPUT_DIR, filename)
        generate_ppt(news_data, output_path)

        json_path = output_path.replace(".pptx", ".json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(news_data, f, ensure_ascii=False, indent=2)
        log.info("Raw data saved to: %s", json_path)

        duration = time.time() - start_time
        record_run(
            success=True,
            ai_count=len(news_data["ai"]),
            auto_count=len(news_data["auto"]),
            top_count=len(news_data["top"]),
            duration=duration,
            ai_items=news_data["ai"],
            auto_items=news_data["auto"],
            top_items=news_data["top"],
            ppt_path=output_path,
            json_path=json_path,
        )

        print(f"\nReport generated successfully!")
        print(f"  PPT:  {output_path}")
        print(f"  JSON: {json_path}")

    except Exception as e:
        duration = time.time() - start_time
        log.error("Task failed: %s", e, exc_info=True)
        record_run(
            success=False,
            ai_count=0, auto_count=0, top_count=0,
            duration=duration,
            error_msg=str(e),
        )
        print(f"\nTask failed: {e}")
        return 1

    return 0


if __name__ == "__main__":
    exit(main())
